import math
from collections import Counter
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


OUTPUT_DIR = Path("outputs")
OUTPUT_DIR.mkdir(exist_ok=True)

TEAMS = [
    ("Corinthians", 53),
    ("Vasco", 57),
    ("Fluminense", 60),
    ("Flamengo", 57),
    ("Internacional", 57),
    ("São Paulo", 57),
    ("Figueirense", 46),
    ("Coritiba", 57),
    ("Botafogo", 52),
    ("Santos", 55),
    ("Palmeiras", 43),
    ("Grêmio", 49),
    ("Atlético-GO", 50),
    ("Bahia", 43),
    ("Atlético-MG", 50),
    ("Cruzeiro", 48),
    ("Athletico-PR", 38),
    ("Ceará", 47),
    ("América-MG", 51),
    ("Avaí", 45),
]


def build_original_dataframe() -> pd.DataFrame:
    df = pd.DataFrame(TEAMS, columns=["Time", "GM"])
    return df.sort_values("Time").reset_index(drop=True)


def build_frequency_table(goals: pd.Series, class_width: int = 5) -> pd.DataFrame:
    min_edge = math.floor(goals.min() / class_width) * class_width
    max_edge = math.ceil(goals.max() / class_width) * class_width + class_width
    bins = np.arange(min_edge, max_edge + class_width, class_width)
    categories = pd.cut(goals, bins=bins, right=False, include_lowest=True)
    freq = categories.value_counts().sort_index()
    cumulative = freq.cumsum()
    rel_freq = freq / freq.sum()
    midpoints = [interval.left + class_width / 2 for interval in freq.index]
    table = pd.DataFrame(
        {
            "Intervalo": [f"{int(interval.left)}-{int(interval.right - 1)}" for interval in freq.index],
            "fi": freq.values,
            "fri": rel_freq.round(3).values,
            "Fac": cumulative.values,
            "xi": midpoints,
        }
    )
    return table


def plot_histogram(goals: pd.Series, class_width: int = 5) -> Path:
    min_edge = math.floor(goals.min() / class_width) * class_width
    max_edge = math.ceil(goals.max() / class_width) * class_width + class_width
    bins = np.arange(min_edge, max_edge + class_width, class_width)
    fig, ax = plt.subplots(figsize=(8, 5))
    ax.hist(goals, bins=bins, edgecolor="black", color="#1f77b4", align="left")
    ax.set_xlabel("Intervalos de gols")
    ax.set_ylabel("Frequência")
    ax.set_title("Histograma de gols marcados (Brasileirão 2011)")
    ax.set_xticks(bins)
    ax.grid(axis="y", alpha=0.3)
    output_path = OUTPUT_DIR / "histograma_gols.png"
    fig.tight_layout()
    fig.savefig(output_path, dpi=300)
    plt.close(fig)
    return output_path


def plot_column_chart(df: pd.DataFrame) -> Path:
    ordered = df.sort_values("GM", ascending=False)
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.bar(ordered["Time"], ordered["GM"], color="#ff7f0e", edgecolor="black")
    ax.set_xlabel("Times")
    ax.set_ylabel("Gols marcados")
    ax.set_title("Gols por time no Brasileirão 2011")
    plt.setp(ax.get_xticklabels(), rotation=45, ha="right")
    ax.grid(axis="y", alpha=0.3)
    output_path = OUTPUT_DIR / "grafico_colunas_gols.png"
    fig.tight_layout()
    fig.savefig(output_path, dpi=300)
    plt.close(fig)
    return output_path


def calculate_statistics(goals: pd.Series) -> dict:
    sorted_goals = goals.sort_values().to_list()
    n = len(sorted_goals)
    mean_val = sum(sorted_goals) / n
    counter = Counter(sorted_goals)
    mode_val = max(counter.items(), key=lambda item: item[1])[0]
    median_val = (sorted_goals[n // 2 - 1] + sorted_goals[n // 2]) / 2 if n % 2 == 0 else sorted_goals[n // 2]
    lower_half = sorted_goals[: n // 2]
    upper_half = sorted_goals[n // 2 :]
    q1 = (lower_half[len(lower_half) // 2 - 1] + lower_half[len(lower_half) // 2]) / 2
    q3 = (upper_half[len(upper_half) // 2 - 1] + upper_half[len(upper_half) // 2]) / 2
    return {
        "mean": mean_val,
        "mode": mode_val,
        "median": median_val,
        "q1": q1,
        "q2": median_val,
        "q3": q3,
    }


def create_presentation(
    df_original: pd.DataFrame,
    frequency_table: pd.DataFrame,
    histogram_path: Path,
    column_chart_path: Path,
    stats: dict,
) -> Path:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = "Gols no Campeonato Brasileiro 2011"
    title_slide.placeholders[1].text = "Análise estatística"

    intro_layout = prs.slide_layouts[1]
    intro_slide = prs.slides.add_slide(intro_layout)
    intro_slide.shapes.title.text = "Introdução"
    intro_body = intro_slide.shapes.placeholders[1].text_frame
    intro_body.text = "Tema: Número de gols marcados no Campeonato Brasileiro de 2011"
    intro_body.add_paragraph().text = "Objetivo: Descrever e interpretar a distribuição de gols por equipe"

    data_slide = prs.slides.add_slide(intro_layout)
    data_slide.shapes.title.text = "Coleta de Dados"
    data_text = data_slide.shapes.placeholders[1].text_frame
    data_text.text = "Fonte: tabela oficial do Brasileirão 2011"
    data_text.add_paragraph().text = "20 observações (uma por equipe)"

    table_slide = prs.slides.add_slide(intro_layout)
    table_slide.shapes.title.text = "Tabela Original"
    rows, cols = df_original.shape
    table = table_slide.shapes.add_table(rows + 1, cols, Inches(0.6), Inches(1.8), Inches(8.8), Inches(4.5)).table
    table.cell(0, 0).text = "Time"
    table.cell(0, 1).text = "Gols"
    for i, (_, row) in enumerate(df_original.iterrows(), start=1):
        table.cell(i, 0).text = row["Time"]
        table.cell(i, 1).text = str(row["GM"])

    freq_slide = prs.slides.add_slide(intro_layout)
    freq_slide.shapes.title.text = "Tabela de Frequências"
    rows, cols = frequency_table.shape
    freq_table_shape = freq_slide.shapes.add_table(
        rows + 1, cols, Inches(0.4), Inches(1.6), Inches(9.2), Inches(4.8)
    ).table
    for j, column in enumerate(frequency_table.columns):
        freq_table_shape.cell(0, j).text = column
    for i, (_, row) in enumerate(frequency_table.iterrows(), start=1):
        for j, value in enumerate(row):
            freq_table_shape.cell(i, j).text = str(value)

    hist_slide = prs.slides.add_slide(prs.slide_layouts[5])
    hist_slide.shapes.title.text = "Histograma"
    hist_slide.shapes.add_picture(str(histogram_path), Inches(1), Inches(1.5), width=Inches(8))

    column_slide = prs.slides.add_slide(prs.slide_layouts[5])
    column_slide.shapes.title.text = "Gráfico de Colunas"
    column_slide.shapes.add_picture(str(column_chart_path), Inches(1), Inches(1.5), width=Inches(8))

    stats_slide = prs.slides.add_slide(intro_layout)
    stats_slide.shapes.title.text = "Medidas Estatísticas"
    stats_frame = stats_slide.shapes.placeholders[1].text_frame
    stats_frame.text = f"Média: {stats['mean']:.2f}"
    stats_frame.add_paragraph().text = f"Moda: {stats['mode']}"
    stats_frame.add_paragraph().text = f"Mediana (Q2): {stats['median']:.2f}"
    stats_frame.add_paragraph().text = f"Q1: {stats['q1']:.2f}"
    stats_frame.add_paragraph().text = f"Q3: {stats['q3']:.2f}"

    concl_slide = prs.slides.add_slide(intro_layout)
    concl_slide.shapes.title.text = "Conclusão"
    concl_frame = concl_slide.shapes.placeholders[1].text_frame
    concl_frame.text = "Os gols concentram-se entre 45 e 57 tentos."
    concl_frame.add_paragraph().text = "Destaca-se a performance do Fluminense (60 gols)."

    final_slide = prs.slides.add_slide(intro_layout)
    final_slide.shapes.title.text = "Encerramento"
    final_slide.shapes.placeholders[1].text = "Obrigado!"

    pptx_path = OUTPUT_DIR / "analise_gols_brasileirao_2011.pptx"
    prs.save(pptx_path)
    return pptx_path


def main() -> None:
    df_original = build_original_dataframe()
    df_original.to_csv(OUTPUT_DIR / "dados_originais.csv", index=False)
    goals = df_original["GM"]
    frequency_table = build_frequency_table(goals)
    frequency_table.to_csv(OUTPUT_DIR / "tabela_frequencias.csv", index=False)
    histogram_path = plot_histogram(goals)
    column_chart_path = plot_column_chart(df_original)
    stats = calculate_statistics(goals)

    pptx_path = create_presentation(df_original, frequency_table, histogram_path, column_chart_path, stats)

    summary = {
        "dados_originais": OUTPUT_DIR / "dados_originais.csv",
        "tabela_frequencias": OUTPUT_DIR / "tabela_frequencias.csv",
        "histograma": histogram_path,
        "grafico_colunas": column_chart_path,
        "apresentacao": pptx_path,
        "estatisticas": stats,
    }
    for key, value in summary.items():
        print(f"{key}: {value}")


if __name__ == "__main__":
    main()
